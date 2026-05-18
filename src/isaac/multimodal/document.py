"""Multimodal Document Understanding — PDF, DOCX, and image OCR.

All backends are lazy-imported so the module loads without optional dependencies.

Supported formats
-----------------
PDF   → pymupdf (local, fast)   → pdfminer.six fallback
DOCX  → python-docx             → plain text extraction
Images → pytesseract + Pillow   → vision LLM fallback
PPTX  → python-pptx             (text extraction)

Usage
-----
    from isaac.multimodal.document import extract_text, extract_pages, analyse_image

    text = extract_text("report.pdf")
    pages = extract_pages("report.pdf")           # list[str], one per page
    result = analyse_image("diagram.png")         # AI-powered image description
"""

from __future__ import annotations

import base64
import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Main entry points
# ---------------------------------------------------------------------------


def extract_text(doc_path: str | Path, max_pages: int = 100) -> str:
    """Extract all text from a document file.

    Dispatches based on file extension: .pdf, .docx, .pptx, or image formats.
    """
    path = Path(doc_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(path, max_pages)
    if suffix in (".docx", ".doc"):
        return _extract_docx(path)
    if suffix in (".pptx", ".ppt"):
        return _extract_pptx(path)
    if suffix in (".txt", ".md", ".rst", ".csv", ".json"):
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".gif"):
        return _ocr_image(path)

    raise ValueError(f"Unsupported document format: {suffix!r}")


def extract_pages(pdf_path: str | Path, max_pages: int = 100) -> list[str]:
    """Extract text page-by-page from a PDF. Returns list[str]."""
    try:
        import fitz  # type: ignore[import-untyped]  # pymupdf

        doc = fitz.open(str(pdf_path))
        pages = []
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            pages.append(page.get_text())
        return pages
    except ImportError:
        logger.warning("pymupdf not installed; falling back to single-chunk extraction")
        text = _extract_pdf(Path(pdf_path), max_pages)
        return [text]


def extract_metadata(doc_path: str | Path) -> dict[str, Any]:
    """Extract document metadata (title, author, page count, etc.)."""
    path = Path(doc_path)
    suffix = path.suffix.lower()

    meta: dict[str, Any] = {"path": str(path), "size_bytes": path.stat().st_size}

    if suffix == ".pdf":
        try:
            import fitz

            doc = fitz.open(str(path))
            raw = doc.metadata or {}
            meta.update(
                {
                    "title": raw.get("title", ""),
                    "author": raw.get("author", ""),
                    "pages": len(doc),
                    "format": "PDF",
                }
            )
        except ImportError:
            pass
    elif suffix in (".docx", ".doc"):
        try:
            import docx  # type: ignore[import-untyped]

            d = docx.Document(str(path))
            props = d.core_properties
            meta.update(
                {
                    "title": props.title or "",
                    "author": props.author or "",
                    "paragraphs": len(d.paragraphs),
                    "format": "DOCX",
                }
            )
        except ImportError:
            pass

    return meta


def analyse_image(
    image_path: str | Path,
    prompt: str = "Describe this image in detail. Extract any text, data, or diagrams.",
    use_llm: bool = True,
) -> dict[str, Any]:
    """Analyse an image using OCR and/or a vision LLM.

    Returns a dict with 'ocr_text', 'description', and 'structured_data'.
    """
    path = Path(image_path)
    result: dict[str, Any] = {
        "path": str(path),
        "ocr_text": "",
        "description": "",
        "structured_data": {},
    }

    # OCR layer
    try:
        result["ocr_text"] = _ocr_image(path)
    except Exception as exc:
        logger.debug("OCR failed: %s", exc)

    # Vision LLM layer
    if use_llm:
        try:
            result["description"] = _vision_llm(path, prompt)
        except Exception as exc:
            logger.debug("Vision LLM failed: %s", exc)

    return result


# ---------------------------------------------------------------------------
# PDF backends
# ---------------------------------------------------------------------------


def _extract_pdf(path: Path, max_pages: int) -> str:
    # 1. pymupdf (fastest, best quality)
    try:
        import fitz

        doc = fitz.open(str(path))
        texts = []
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            texts.append(page.get_text())
        return "\n\n".join(texts)
    except ImportError:
        pass

    # 2. pdfminer.six fallback
    try:
        from pdfminer.high_level import extract_text as pm_extract  # type: ignore[import-untyped]

        return pm_extract(str(path), maxpages=max_pages)
    except ImportError as exc:
        raise ImportError(
            "PDF extraction requires pymupdf or pdfminer.six. Install with: pip install pymupdf"
        ) from exc


# ---------------------------------------------------------------------------
# DOCX backend
# ---------------------------------------------------------------------------


def _extract_docx(path: Path) -> str:
    try:
        import docx  # type: ignore[import-untyped]

        d = docx.Document(str(path))
        return "\n".join(para.text for para in d.paragraphs if para.text.strip())
    except ImportError as exc:
        raise ImportError(
            "DOCX extraction requires python-docx. Install with: pip install python-docx"
        ) from exc


# ---------------------------------------------------------------------------
# PPTX backend
# ---------------------------------------------------------------------------


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation(str(path))
        slides_text = []
        for slide in prs.slides:
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
            slides_text.append("\n".join(parts))
        return "\n\n---\n\n".join(slides_text)
    except ImportError as exc:
        raise ImportError(
            "PPTX extraction requires python-pptx. Install with: pip install python-pptx"
        ) from exc


# ---------------------------------------------------------------------------
# OCR backend
# ---------------------------------------------------------------------------


def _ocr_image(path: Path) -> str:
    try:
        import pytesseract  # type: ignore[import-untyped]
        from PIL import Image  # type: ignore[import-untyped]

        img = Image.open(str(path))
        return pytesseract.image_to_string(img)
    except ImportError as exc:
        raise ImportError(
            "Image OCR requires pytesseract and Pillow. "
            "Install with: pip install pytesseract Pillow\n"
            "Also install Tesseract OCR: https://github.com/UB-Mannheim/tesseract/wiki"
        ) from exc


# ---------------------------------------------------------------------------
# Vision LLM backend
# ---------------------------------------------------------------------------


def _vision_llm(path: Path, prompt: str) -> str:
    """Send image to a vision-capable LLM and return the response."""
    from langchain_core.messages import HumanMessage

    from isaac.llm.provider import get_llm

    img_b64 = base64.b64encode(path.read_bytes()).decode()
    suffix = path.suffix.lower().lstrip(".")
    mime = {
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "png": "image/png",
        "webp": "image/webp",
        "gif": "image/gif",
    }.get(suffix, "image/png")

    llm = get_llm("strong")
    msg = HumanMessage(
        content=[
            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{img_b64}"}},
            {"type": "text", "text": prompt},
        ]
    )
    response = llm.invoke([msg])
    return str(response.content)
